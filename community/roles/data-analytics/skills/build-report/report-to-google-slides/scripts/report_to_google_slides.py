#!/usr/bin/env python3
"""Convert an HTML analytics report into an editable PPTX deck.

The script does not call Google APIs. It builds a local PPTX plus deterministic
planning and preflight artifacts that an agent can import into native Google
Slides with the Google Drive connector.
"""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

try:
    from bs4 import BeautifulSoup, NavigableString, Tag
except ImportError as exc:  # pragma: no cover - dependency failure path
    raise SystemExit(
        "This helper requires beautifulsoup4. Install it in the active Python "
        "environment before running report_to_google_slides.py."
    ) from exc

try:
    from PIL import Image, ImageChops, ImageFont
except ImportError as exc:  # pragma: no cover - dependency failure path
    raise SystemExit(
        "This helper requires pillow. Install it in the active Python "
        "environment before running report_to_google_slides.py."
    ) from exc

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency failure path
    raise SystemExit(
        "This helper requires python-pptx. Install it in the active Python "
        "environment before running report_to_google_slides.py."
    ) from exc


SLIDE_W = 13.333333
SLIDE_H = 7.5
PX_PER_PT = 96 / 72
DEFAULT_COLORS = {
    "surface": "#FCFCFD",
    "ink": "#1F2430",
    "muted": "#6F768A",
    "grid": "#E6E8F0",
    "axis": "#D7DBE7",
    "blue": "#2E4780",
    "blue_light": "#CEDFFE",
    "gold": "#B8A037",
    "neutral": "#C5CAD3",
    "neutral_dark": "#464C55",
}

EXECUTIVE_SECTION_CANDIDATES = ["Executive Summary", "Summary", "Overview"]
CONCLUSION_SECTION_CANDIDATES = ["Conclusions", "Conclusion", "Key Takeaways", "Takeaways"]
IMPLICATION_SECTION_CANDIDATES = [
    "Implications",
    "Recommendations",
    "Recommended Actions",
    "Next Steps",
]
CAVEAT_SECTION_CANDIDATES = [
    "Caveats And Assumptions",
    "Caveats",
    "Assumptions",
    "Limitations",
    "Methodology Notes",
]
SOURCE_SECTION_CANDIDATES = [
    "Sources And Reproducibility",
    "Sources",
    "Reproducibility",
    "Methodology",
    "Appendix",
]
BLOCKED_SVG_TAGS = {
    "audio",
    "base",
    "canvas",
    "embed",
    "foreignobject",
    "iframe",
    "link",
    "meta",
    "object",
    "script",
    "video",
}


@dataclass
class MetricCard:
    label: str
    value: str
    note: str = ""


@dataclass
class ChartItem:
    index: int
    title: str
    takeaway: str
    caveat: str
    svg_markup: str
    css_text: str = ""
    image_path: str = ""
    width_px: int = 0
    height_px: int = 0


@dataclass
class TableItem:
    headers: list[str]
    rows: list[list[str]]
    title: str = ""
    note: str = ""


@dataclass
class LinkItem:
    text: str
    href: str


@dataclass
class Section:
    title: str
    paragraphs: list[str] = field(default_factory=list)
    bullets: list[str] = field(default_factory=list)
    links: list[LinkItem] = field(default_factory=list)


@dataclass
class Report:
    source: str
    title: str
    eyebrow: str
    lede: str
    metrics: list[MetricCard]
    sections: dict[str, Section]
    charts: list[ChartItem]
    tables: list[TableItem]
    colors: dict[str, str]


@dataclass
class Check:
    name: str
    status: str
    message: str
    details: dict[str, Any] = field(default_factory=dict)


@dataclass
class Preflight:
    checks: list[Check] = field(default_factory=list)

    def add(self, name: str, ok: bool, message: str, details: dict[str, Any] | None = None) -> None:
        self.checks.append(
            Check(
                name=name,
                status="passed" if ok else "failed",
                message=message,
                details=details or {},
            )
        )

    @property
    def passed(self) -> bool:
        return all(check.status == "passed" for check in self.checks)

    def to_json(self) -> dict[str, Any]:
        failed = [check for check in self.checks if check.status != "passed"]
        return {
            "status": "passed" if not failed else "failed",
            "summary": {
                "checks": len(self.checks),
                "passed": len(self.checks) - len(failed),
                "failed": len(failed),
            },
            "checks": [asdict(check) for check in self.checks],
        }


def collapse_ws(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def first_sentence(text: str, max_chars: int = 140) -> str:
    text = collapse_ws(text)
    if not text:
        return ""
    match = re.search(r"(?<=[.!?])\s+", text)
    sentence = text[: match.start()].strip() if match else text
    if len(sentence) <= max_chars:
        return sentence
    return sentence[: max_chars - 3].rsplit(" ", 1)[0] + "..."


def normalized_label(text: str) -> str:
    text = collapse_ws(text).lower().replace("&", " and ")
    return re.sub(r"[^a-z0-9]+", " ", text).strip()


def find_section(report: Report, candidates: list[str]) -> Section | None:
    by_label = {normalized_label(title): section for title, section in report.sections.items()}
    for candidate in candidates:
        section = by_label.get(normalized_label(candidate))
        if section:
            return section

    candidate_words = [set(normalized_label(candidate).split()) for candidate in candidates]
    for title, section in report.sections.items():
        title_words = set(normalized_label(title).split())
        if any(words and words.issubset(title_words) for words in candidate_words):
            return section
    return None


def section_items_by_candidates(report: Report, candidates: list[str]) -> list[str]:
    section = find_section(report, candidates)
    if not section:
        return []
    return section.paragraphs + section.bullets


def section_title_by_candidates(report: Report, candidates: list[str], fallback: str) -> str:
    section = find_section(report, candidates)
    return section.title if section and section.title else fallback


def has_section_content(report: Report, candidates: list[str]) -> bool:
    section = find_section(report, candidates)
    return bool(section and (section.paragraphs or section.bullets or section.links))


def resolve_source(source: str) -> Path:
    candidate = Path(source).expanduser()
    if candidate.is_file():
        return candidate.resolve()

    if source.startswith(("http://", "https://")):
        raise SystemExit(
            "This helper only accepts a local HTML file path. Retrieve the report "
            "using the user's environment-specific instructions, then rerun with "
            "the local file path."
        )

    raise SystemExit(f"Source does not exist or is not a file: {source}")


def read_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    if len(raw) < 1000:
        raise SystemExit(f"HTML file is too small to be a report: {path}")
    return raw


def validate_report_html(soup: BeautifulSoup, path: Path) -> None:
    title = collapse_ws(soup.title.get_text(" ", strip=True) if soup.title else "")
    body_text = collapse_ws(soup.get_text(" ", strip=True))
    has_report_shape = bool(soup.find("h1") and soup.find("h2")) or bool(
        soup.find_all(["svg", "table"])
    )
    auth_markers = [
        "sign in to your account",
        "login.microsoftonline.com",
        "convergedsignin",
        "jsdisabled",
    ]
    lowered = f"{title} {body_text[:1000]}".lower()
    if any(marker in lowered for marker in auth_markers) or not has_report_shape:
        raise SystemExit(
            f"{path} does not look like the report HTML. It appears to be an auth "
            "or non-report page; provide the actual local HTML report file and rerun."
        )


def extract_report_css(soup: BeautifulSoup) -> str:
    return "\n".join(style.get_text() for style in soup.find_all("style"))


def parse_css_vars(soup: BeautifulSoup, style_text: str | None = None) -> dict[str, str]:
    colors = DEFAULT_COLORS.copy()
    style_text = style_text if style_text is not None else extract_report_css(soup)
    for name, value in re.findall(r"--([a-zA-Z0-9_-]+)\s*:\s*(#[0-9a-fA-F]{3,6})", style_text):
        key = name.replace("-", "_")
        if key in colors or key in {
            "surface",
            "ink",
            "muted",
            "grid",
            "axis",
            "blue",
            "gold",
            "neutral",
        }:
            colors[key] = normalize_hex(value)
    return colors


def normalize_hex(value: str) -> str:
    value = value.strip()
    if re.fullmatch(r"#[0-9a-fA-F]{3}", value):
        return "#" + "".join(ch * 2 for ch in value[1:]).upper()
    if re.fullmatch(r"#[0-9a-fA-F]{6}", value):
        return value.upper()
    return "#000000"


def text_of(node: Tag | None) -> str:
    if not node:
        return ""
    return collapse_ws(node.get_text(" ", strip=True))


def links_of(node: Tag) -> list[LinkItem]:
    links: list[LinkItem] = []
    for anchor in node.find_all("a", href=True):
        href = collapse_ws(str(anchor.get("href") or ""))
        if not href or href.startswith("#"):
            continue
        label = text_of(anchor) or href
        links.append(LinkItem(text=label, href=href))
    return links


def add_section_links(section: Section, links: list[LinkItem]) -> None:
    seen = {(link.text, link.href) for link in section.links}
    for link in links:
        key = (link.text, link.href)
        if key not in seen:
            section.links.append(link)
            seen.add(key)


def parse_metrics(soup: BeautifulSoup) -> list[MetricCard]:
    metrics: list[MetricCard] = []
    for metric in soup.select(".metric"):
        label = text_of(metric.select_one(".label")) or text_of(metric.find("span"))
        value = text_of(metric.select_one(".value")) or text_of(metric.find("strong"))
        note = text_of(metric.select_one(".note"))
        if not note:
            value_node = metric.select_one(".value") or metric.find("strong")
            label_node = metric.select_one(".label") or metric.find("span")
            excluded = {id(node) for node in (value_node, label_node) if node}
            note_parts = [
                collapse_ws(str(part))
                for part in metric.stripped_strings
                if collapse_ws(str(part))
                and collapse_ws(str(part)) not in {label, value}
                and id(getattr(part, "parent", None)) not in excluded
            ]
            note = " ".join(note_parts)
        metrics.append(
            MetricCard(
                label=label,
                value=value,
                note=note,
            )
        )
    return [m for m in metrics if m.label or m.value]


def parse_table(table: Tag) -> TableItem:
    parsed_rows: list[tuple[list[str], bool]] = []
    for tr in table.find_all("tr"):
        cell_nodes = tr.find_all(["th", "td"])
        cells = [text_of(cell) for cell in cell_nodes]
        if cells:
            parsed_rows.append((cells, any(cell.name == "th" for cell in cell_nodes)))
    if not parsed_rows:
        return TableItem(headers=[], rows=[])
    width = max(len(row) for row, _ in parsed_rows)
    header_idx = next((idx for idx, (_, has_th) in enumerate(parsed_rows) if has_th), None)
    if header_idx is None:
        headers = [f"Column {idx}" for idx in range(1, width + 1)]
        rows = [row + [""] * (width - len(row)) for row, _ in parsed_rows]
        return TableItem(headers=headers, rows=rows)
    headers = parsed_rows[header_idx][0]
    headers = headers + [f"Column {idx}" for idx in range(len(headers) + 1, width + 1)]
    rows = [
        row + [""] * (width - len(row))
        for idx, (row, _) in enumerate(parsed_rows)
        if idx != header_idx
    ]
    return TableItem(headers=headers, rows=rows)


def table_context_title(table: Tag, fallback: str = "") -> str:
    heading = table.find_previous(["h4", "h3", "h2"])
    return text_of(heading) or fallback


def table_context_note(table: Tag) -> str:
    for sibling in table.next_siblings:
        if isinstance(sibling, NavigableString):
            continue
        if not isinstance(sibling, Tag):
            continue
        if sibling.name in {"h2", "h3", "h4", "table"}:
            return ""
        if sibling.name in {"p", "figcaption"} or "caption" in set(sibling.get("class") or []):
            return text_of(sibling)
    return ""


def find_matching_table(tables: list[TableItem], parsed: TableItem) -> TableItem | None:
    for table in tables:
        if table.headers == parsed.headers and table.rows == parsed.rows:
            return table
    return None


def parse_svg_length(value: Any, fallback: float) -> float:
    match = re.search(r"\d+(?:\.\d+)?", str(value or ""))
    if not match:
        return fallback
    try:
        return float(match.group(0))
    except ValueError:
        return fallback


def chart_dimensions(svg: Tag) -> tuple[int, int]:
    view_box = svg.get("viewBox") or svg.get("viewbox")
    if view_box:
        raw_parts = [part for part in re.split(r"[,\s]+", str(view_box).strip()) if part]
        try:
            parts = [float(part) for part in raw_parts]
        except ValueError:
            parts = []
        if len(parts) == 4 and parts[2] > 0 and parts[3] > 0:
            return int(math.ceil(parts[2])), int(math.ceil(parts[3]))
    width = parse_svg_length(svg.get("width"), 900)
    height = parse_svg_length(svg.get("height"), 520)
    return int(math.ceil(width)), int(math.ceil(height))


def looks_like_chart_title_candidate(text: str) -> bool:
    text = collapse_ws(text)
    if not 8 <= len(text) <= 180:
        return False
    if re.fullmatch(r"[a-z]+/[a-z0-9.+-]+", text.lower()):
        return False
    lowered = normalized_label(text)
    if "matplotlib" in lowered or "http" in text.lower():
        return False
    if len(re.findall(r"[A-Za-z]", text)) < 4:
        return False
    words = re.findall(r"[A-Za-z][A-Za-z0-9_-]*", text)
    if len(words) < 2:
        return False
    numeric_ratio = len(re.findall(r"[0-9]", text)) / max(1, len(text))
    if numeric_ratio > 0.45:
        return False
    axis_or_legend_labels = {
        "source",
        "notes",
        "legend",
        "date",
        "day",
        "week",
        "month",
        "year",
        "image svg xml",
    }
    return lowered not in axis_or_legend_labels


def chart_text_candidates(svg: Tag) -> list[str]:
    candidates: list[str] = []
    for value in svg.stripped_strings:
        text = collapse_ws(str(value))
        if looks_like_chart_title_candidate(text):
            candidates.append(text)

    deduped: list[str] = []
    seen: set[str] = set()
    for candidate in candidates:
        key = normalized_label(candidate)
        if key and key not in seen:
            deduped.append(candidate)
            seen.add(key)
    return deduped


def chart_semantic_label(svg: Tag) -> str:
    for attr in ("aria-label", "data-title", "title"):
        value = collapse_ws(str(svg.get(attr) or ""))
        if value and looks_like_chart_title_candidate(value):
            return value
    for tag_name in ("title", "desc"):
        for node in svg.find_all(tag_name):
            value = text_of(node)
            if value and looks_like_chart_title_candidate(value):
                return value
    return ""


def chart_title_score(text: str) -> float:
    words = re.findall(r"[A-Za-z][A-Za-z0-9_-]*", text)
    score = min(len(words), 14) * 0.8
    score += min(len(text), 120) / 45
    if re.search(r"[.!?]$", text.strip()):
        score += 2.0
    if "," in text or ":" in text or ";" in text:
        score += 0.8
    if len(words) <= 4 and not re.search(r"[.!?]$", text.strip()):
        score -= 4.0
    if len(text) < 25:
        score -= 2.0
    return score


def parse_chart_title(svg: Tag) -> str:
    semantic_label = chart_semantic_label(svg)
    if semantic_label:
        return first_sentence(semantic_label, 120)
    candidates = chart_text_candidates(svg)
    if candidates:
        return first_sentence(max(candidates, key=chart_title_score), 120)
    text = collapse_ws(svg.get_text(" ", strip=True))
    return first_sentence(text, 100) or "Source chart from report"


def is_material_chart_svg(svg: Tag, width_px: int, height_px: int) -> bool:
    if width_px < 180 or height_px < 110:
        return False
    if chart_semantic_label(svg) or chart_text_candidates(svg):
        return True
    visual_marks = svg.find_all(
        ["circle", "ellipse", "line", "path", "polygon", "polyline", "rect"]
    )
    if width_px >= 280 and height_px >= 170 and len(visual_marks) >= 3:
        return True
    return width_px >= 420 and height_px >= 240 and bool(visual_marks)


def is_chart_note_text(text: str) -> bool:
    normalized = normalized_label(text)
    note_prefixes = [
        "note",
        "notes",
        "source",
        "definition",
        "current window",
        "previous window",
        "bars are",
        "bars show",
        "values are",
    ]
    return any(normalized.startswith(prefix) for prefix in note_prefixes)


def chart_takeaway(chart_title: str, next_para: str) -> str:
    candidate = first_sentence(next_para, 145)
    if not candidate or is_chart_note_text(candidate):
        return chart_title
    if normalized_label(candidate).startswith(("implication", "recommendation")):
        return chart_title
    return candidate


def chart_context_text(
    section_title: str,
    chart_title: str,
    next_para: str,
    prior_items: list[str],
) -> str:
    candidates: list[str] = []
    if next_para:
        candidates.append(next_para)
    candidates.extend(reversed(prior_items))
    chart_key = normalize_sentence_key(chart_title)

    for allow_note_text in (False, True):
        for candidate in candidates:
            raw_candidate = collapse_ws(candidate)
            sentence = first_sentence(raw_candidate, 230)
            sentence_key = normalize_sentence_key(sentence)
            if not sentence or sentence_key == chart_key:
                continue
            if not allow_note_text and is_chart_note_text(sentence):
                continue
            return shorten_for_bullet(raw_candidate, 360)
    return f"This chart supports the report section: {section_title}."


def iter_content_blocks(node: Tag) -> list[Tag]:
    if node.name in {"h2", "h3", "h4"}:
        return []
    classes = set(node.get("class") or [])
    if node.name in {"p", "ul", "ol", "table", "svg"}:
        return [node]
    if "figure" in classes and node.find("svg"):
        return [node]
    direct_svg = node.find("svg", recursive=False)
    direct_text_blocks = node.find(["p", "ul", "ol", "table"], recursive=False)
    if direct_svg and not direct_text_blocks:
        return [node]

    blocks: list[Tag] = []
    for child in node.children:
        if isinstance(child, Tag):
            blocks.extend(iter_content_blocks(child))
    return blocks


def heading_level(node: Tag) -> int:
    if node.name and re.fullmatch(r"h[1-6]", node.name):
        return int(node.name[1])
    return 7


def section_content_blocks(heading: Tag) -> list[Tag]:
    blocks: list[Tag] = []
    current_level = heading_level(heading)
    for sibling in heading.next_siblings:
        if isinstance(sibling, NavigableString):
            continue
        if not isinstance(sibling, Tag):
            continue
        if re.fullmatch(r"h[1-6]", sibling.name or "") and heading_level(sibling) <= current_level:
            break
        blocks.extend(iter_content_blocks(sibling))
    return blocks


def context_text_from_block(block: Tag) -> str:
    if block.name == "p":
        return text_of(block)
    if block.name in {"ul", "ol"}:
        return " ".join(text_of(li) for li in block.find_all("li", recursive=False))
    return ""


def following_context(blocks: list[Tag], start_idx: int) -> str:
    for block in blocks[start_idx + 1 :]:
        text = context_text_from_block(block)
        if text:
            return text
    return ""


def nearest_heading_title(node: Tag, fallback: str) -> str:
    heading = node.find_previous(["h4", "h3", "h2", "h1"])
    heading_text = text_of(heading) if isinstance(heading, Tag) else ""
    return heading_text or fallback


def nearby_prior_items(node: Tag, limit: int = 4) -> list[str]:
    items: list[str] = []
    for prev in node.find_all_previous(["p", "li"], limit=limit):
        text = text_of(prev)
        if text:
            items.append(text)
    return list(reversed(items))


def nearby_next_context(node: Tag, limit: int = 4) -> str:
    for nxt in node.find_all_next(["p", "li"], limit=limit):
        text = text_of(nxt)
        if text:
            return text
    return ""


def add_chart_item(
    charts: list[ChartItem],
    svg_container: Tag,
    section_title: str,
    next_para: str,
    prior_items: list[str],
    style_text: str,
    parsed_svg_ids: set[int],
) -> None:
    svg = svg_container if svg_container.name == "svg" else svg_container.find("svg")
    if not isinstance(svg, Tag) or id(svg) in parsed_svg_ids:
        return
    width_px, height_px = chart_dimensions(svg)
    if not is_material_chart_svg(svg, width_px, height_px):
        parsed_svg_ids.add(id(svg))
        return
    parsed_svg_ids.add(id(svg))
    chart_title = parse_chart_title(svg)
    if chart_title == "Source chart from report":
        chart_title = section_title
    context_text = chart_context_text(section_title, chart_title, next_para, prior_items)
    charts.append(
        ChartItem(
            index=len(charts) + 1,
            title=chart_title,
            takeaway=chart_takeaway(chart_title, context_text),
            caveat=context_text,
            svg_markup=str(svg),
            css_text=style_text,
            width_px=width_px,
            height_px=height_px,
        )
    )


def parse_sections_and_charts(
    soup: BeautifulSoup,
    style_text: str,
) -> tuple[dict[str, Section], list[ChartItem], list[TableItem]]:
    sections: dict[str, Section] = {}
    charts: list[ChartItem] = []
    tables: list[TableItem] = []
    parsed_svg_ids: set[int] = set()
    parsed_table_ids: set[int] = set()

    for heading in soup.find_all(["h2", "h3"]):
        title = text_of(heading)
        if not title:
            continue
        section = Section(title=title)
        blocks = section_content_blocks(heading)
        for block_idx, block in enumerate(blocks):
            if block.name == "p":
                paragraph = text_of(block)
                if paragraph:
                    section.paragraphs.append(paragraph)
                add_section_links(section, links_of(block))
            elif block.name in {"ul", "ol"}:
                for li in block.find_all("li", recursive=False):
                    item = text_of(li)
                    if item:
                        section.bullets.append(item)
                    add_section_links(section, links_of(li))
            elif block.name == "table":
                parsed_table_ids.add(id(block))
                table = parse_table(block)
                if table.headers or table.rows:
                    table.title = table_context_title(block, title)
                    table.note = table_context_note(block)
                    tables.append(table)
            elif block.name == "svg" or block.find("svg"):
                add_chart_item(
                    charts,
                    block,
                    title,
                    following_context(blocks, block_idx),
                    section.paragraphs + section.bullets,
                    style_text,
                    parsed_svg_ids,
                )
        sections[title] = section

    fallback_title = text_of(soup.find("h1")) or text_of(soup.title) or "Chart evidence"
    for svg in soup.find_all("svg"):
        if id(svg) in parsed_svg_ids:
            continue
        add_chart_item(
            charts,
            svg,
            nearest_heading_title(svg, fallback_title),
            nearby_next_context(svg),
            nearby_prior_items(svg),
            style_text,
            parsed_svg_ids,
        )

    for table in soup.find_all("table"):
        if id(table) in parsed_table_ids:
            continue
        parsed = parse_table(table)
        if not (parsed.headers or parsed.rows):
            continue
        parsed.title = table_context_title(table)
        parsed.note = table_context_note(table)
        match = find_matching_table(tables, parsed)
        if match:
            if parsed.title and not match.title:
                match.title = parsed.title
            if parsed.note and not match.note:
                match.note = parsed.note
        else:
            tables.append(parsed)

    return sections, charts, tables


def parse_report(path: Path) -> Report:
    raw = read_html(path)
    soup = BeautifulSoup(raw, "html.parser")
    validate_report_html(soup, path)
    style_text = extract_report_css(soup)
    colors = parse_css_vars(soup, style_text)
    title = text_of(soup.find("h1")) or text_of(soup.title) or "Analytics Report"
    eyebrow = text_of(soup.select_one(".eyebrow"))
    lede = text_of(soup.select_one(".lede"))
    metrics = parse_metrics(soup)
    sections, charts, tables = parse_sections_and_charts(soup, style_text)
    if not lede:
        executive_labels = {normalized_label(label) for label in EXECUTIVE_SECTION_CANDIDATES}
        for section_title, section in sections.items():
            if normalized_label(section_title) in executive_labels and section.paragraphs:
                lede = first_sentence(section.paragraphs[0], 220)
                break
    return Report(
        source=str(path),
        title=title,
        eyebrow=eyebrow,
        lede=lede,
        metrics=metrics,
        sections=sections,
        charts=charts,
        tables=tables,
        colors=colors,
    )


def write_json(path: Path, data: Any) -> None:
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def rgb(value: str) -> RGBColor:
    value = normalize_hex(value)
    return RGBColor(int(value[1:3], 16), int(value[3:5], 16), int(value[5:7], 16))


def find_font_path(bold: bool = False) -> str | None:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
        if bold
        else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
        if bold
        else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def pil_font(
    font_size_pt: float, bold: bool = False
) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    path = find_font_path(bold)
    size_px = max(8, round(font_size_pt * PX_PER_PT))
    if path:
        return ImageFont.truetype(path, size_px)
    return ImageFont.load_default()


def line_width_px(text: str, font: ImageFont.FreeTypeFont | ImageFont.ImageFont) -> float:
    bbox = font.getbbox(text)
    return max(0, bbox[2] - bbox[0])


def wrap_lines(
    text: str, font_size_pt: float, max_width_pt: float, bold: bool = False
) -> list[str]:
    font = pil_font(font_size_pt, bold=bold)
    max_width_px = max_width_pt * PX_PER_PT
    lines: list[str] = []
    for paragraph in (text or "").splitlines() or [""]:
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        line = words[0]
        for word in words[1:]:
            candidate = f"{line} {word}"
            if line_width_px(candidate, font) <= max_width_px:
                line = candidate
            else:
                lines.append(line)
                line = word
        lines.append(line)
    return lines


def measured_height_pt(
    text: str,
    font_size_pt: float,
    max_width_pt: float,
    bold: bool = False,
    line_spacing: float = 1.18,
    paragraph_gap_pt: float = 2.5,
) -> float:
    paragraphs = (text or "").splitlines() or [""]
    total = 0.0
    for paragraph in paragraphs:
        lines = wrap_lines(paragraph, font_size_pt, max_width_pt, bold=bold)
        total += max(1, len(lines)) * font_size_pt * line_spacing + paragraph_gap_pt
    return total


def best_font_size(
    text: str,
    max_width_pt: float,
    max_height_pt: float,
    start_pt: float,
    min_pt: float,
    bold: bool = False,
) -> tuple[float, bool]:
    size = start_pt
    while size >= min_pt:
        if measured_height_pt(text, size, max_width_pt, bold=bold) <= max_height_pt:
            return size, True
        size -= 1
    return min_pt, measured_height_pt(text, min_pt, max_width_pt, bold=bold) <= max_height_pt


def set_shape_fill(shape: Any, color: str, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    if transparency:
        shape.fill.transparency = transparency


def set_shape_line(
    shape: Any, color: str, width_pt: float = 1.0, transparency: float = 0.0
) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width_pt)
    if transparency:
        shape.line.transparency = transparency


def add_box(
    slide: Any,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str | None = None,
    radius: bool = False,
) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    set_shape_fill(shape, fill)
    set_shape_line(shape, line or fill, width_pt=0.75)
    return shape


def add_text(
    slide: Any,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_size: float,
    color: str,
    preflight: Preflight,
    *,
    bold: bool = False,
    min_font: float = 8,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    text = collapse_ws(text)
    margin_pt = 4
    chosen, fits = best_font_size(
        text,
        max_width_pt=w * 72 - margin_pt * 2,
        max_height_pt=h * 72 - margin_pt * 2,
        start_pt=font_size,
        min_pt=min_font,
        bold=bold,
    )
    preflight.add(
        f"text_fit:{name}",
        fits,
        f"{name} text fits at {chosen:.0f}pt" if fits else f"{name} may overflow at minimum size",
        {"font_size_pt": chosen, "box": [x, y, w, h], "text": text[:180]},
    )
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(margin_pt)
    tf.margin_right = Pt(margin_pt)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = "Arial"
    run.font.size = Pt(chosen)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def apply_native_bullet(paragraph: Any) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag.endswith(("}buChar", "}buAutoNum", "}buBlip", "}buNone")):
            p_pr.remove(child)
    p_pr.set("marL", "285750")
    p_pr.set("indent", "-142875")
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "\u2022")
    p_pr.insert(0, bullet)


def add_bullets(
    slide: Any,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: list[str],
    font_size: float,
    color: str,
    preflight: Preflight,
    *,
    max_items: int = 5,
    min_font: float = 9,
    links: list[str | None] | None = None,
    link_color: str | None = None,
    space_after_pt: float = 5.5,
) -> Any:
    clean = [collapse_ws(item) for item in bullets if collapse_ws(item)][:max_items]
    clean_links: list[str | None] = []
    if links:
        clean_links = [links[idx] if idx < len(links) else None for idx in range(len(clean))]
    else:
        clean_links = [None] * len(clean)
    bullet_text = "\n".join(clean)
    chosen, fits = best_font_size(
        bullet_text,
        max_width_pt=w * 72 - 24,
        max_height_pt=h * 72 - 8,
        start_pt=font_size,
        min_pt=min_font,
    )
    preflight.add(
        f"text_fit:{name}",
        fits,
        f"{name} bullets fit at {chosen:.0f}pt" if fits else f"{name} bullets may overflow",
        {"items": len(clean), "font_size_pt": chosen, "box": [x, y, w, h]},
    )
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for idx, (item, href) in enumerate(zip(clean, clean_links, strict=False)):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(space_after_pt)
        p.level = 0
        apply_native_bullet(p)
        run = p.add_run()
        run.text = item
        run.font.name = "Arial"
        run.font.size = Pt(chosen)
        run.font.color.rgb = rgb(link_color if href and link_color else color)
        if href:
            run.hyperlink.address = href
    return shape


def add_source(slide: Any, text: str, colors: dict[str, str], preflight: Preflight) -> None:
    add_text(
        slide,
        "source-note",
        0.62,
        7.02,
        12.1,
        0.34,
        text,
        7.5,
        colors.get("muted", DEFAULT_COLORS["muted"]),
        preflight,
        min_font=6,
    )


def blank_slide(prs: Presentation, colors: dict[str, str]) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(
        slide,
        "slide-background",
        0,
        0,
        SLIDE_W,
        SLIDE_H,
        colors.get("surface", DEFAULT_COLORS["surface"]),
    )
    return slide


def sanitize_css(css_text: str) -> str:
    css_text = re.sub(r"@import[^;]+;", "", css_text, flags=re.IGNORECASE)
    css_text = re.sub(r"@font-face\s*\{[^}]*\}", "", css_text, flags=re.IGNORECASE | re.DOTALL)
    css_text = re.sub(r"expression\s*\([^)]*\)", "none", css_text, flags=re.IGNORECASE)

    def replace_url(match: re.Match[str]) -> str:
        target = match.group(1).strip().strip("\"'")
        return match.group(0) if target.startswith("#") else "none"

    css_text = re.sub(r"url\(\s*([^)]*?)\s*\)", replace_url, css_text, flags=re.IGNORECASE)
    return css_text.replace("</style", "<\\/style")


def sanitize_svg_markup(svg_markup: str) -> str:
    soup = BeautifulSoup(svg_markup, "html.parser")
    for tag in soup.find_all(BLOCKED_SVG_TAGS):
        tag.decompose()
    for tag in soup.find_all(True):
        if tag.name == "style":
            tag.string = sanitize_css(tag.get_text())
        for attr in list(tag.attrs):
            attr_lower = attr.lower()
            value = tag.attrs[attr]
            value_text = " ".join(value) if isinstance(value, list) else str(value)
            if attr_lower.startswith("on"):
                del tag.attrs[attr]
            elif attr_lower in {"href", "xlink:href", "src"}:
                target = collapse_ws(value_text)
                if target and not target.startswith("#"):
                    del tag.attrs[attr]
            elif attr_lower == "style":
                tag.attrs[attr] = sanitize_css(value_text)
            elif "url(" in value_text.lower():
                tag.attrs[attr] = sanitize_css(value_text)
    svg = soup.find("svg")
    return str(svg) if svg else ""


def svg_with_embedded_css(svg_markup: str, css_text: str) -> str:
    css_text = css_text.strip()
    if not css_text:
        return svg_markup
    soup = BeautifulSoup(svg_markup, "html.parser")
    svg = soup.find("svg")
    if svg is None:
        return svg_markup
    style = soup.new_tag("style")
    style.string = css_text
    svg.insert(0, style)
    return str(svg)


def render_svg_chart(chart: ChartItem, out_dir: Path) -> tuple[Path | None, dict[str, Any]]:
    asset_dir = out_dir / "assets" / "charts"
    asset_dir.mkdir(parents=True, exist_ok=True)
    image_path = asset_dir / f"chart_{chart.index:02d}.png"
    safe_svg = sanitize_svg_markup(chart.svg_markup)
    if not safe_svg:
        return None, {"error": "chart SVG was empty after sanitization"}
    report_css = sanitize_css(chart.css_text)
    html = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        "<style>"
        "html,body{margin:0;padding:0;background:white;}svg{display:block;}"
        f"{report_css}"
        "</style>"
        "</head><body>" + safe_svg + "</body></html>"
    )
    try:
        import cairosvg  # type: ignore[import-not-found]

        cairosvg.svg2png(
            bytestring=svg_with_embedded_css(safe_svg, report_css).encode("utf-8"),
            write_to=str(image_path),
        )
        return image_path, {"sanitized": True, "renderer": "cairosvg"}
    except Exception as cairosvg_error:
        cairo_details = {"cairosvg_error": str(cairosvg_error)[:500]}

    try:
        doc_skill_scripts = Path(__file__).resolve().parents[2] / "report-to-google-doc" / "scripts"
        if str(doc_skill_scripts) not in sys.path:
            sys.path.insert(0, str(doc_skill_scripts))
        from report_to_google_doc.rendering import render_simple_svg_with_pillow

        render_simple_svg_with_pillow(
            svg_with_embedded_css(safe_svg, report_css),
            image_path,
            max(360, chart.width_px),
            max(220, chart.height_px),
        )
        return image_path, {"sanitized": True, "renderer": "pillow_svg_subset", **cairo_details}
    except Exception as pillow_error:
        pillow_details = {"pillow_svg_error": str(pillow_error)[:500]}

    try:
        from playwright.sync_api import sync_playwright

        with sync_playwright() as playwright:
            browser = playwright.chromium.launch(headless=True)
            context = browser.new_context(
                viewport={
                    "width": max(360, min(1600, chart.width_px + 8)),
                    "height": max(220, min(1100, chart.height_px + 8)),
                    "device_scale_factor": 2,
                },
                java_script_enabled=False,
            )
            page = context.new_page()
            page.route("**/*", lambda route: route.abort())
            page.set_content(html)
            page.locator("svg").screenshot(path=str(image_path))
            context.close()
            browser.close()
    except Exception as exc:
        svg_path = asset_dir / f"chart_{chart.index:02d}.svg"
        svg_path.write_text(safe_svg, encoding="utf-8")
        return None, {
            "error": str(exc)[:500],
            "sanitized_svg": str(svg_path),
            **cairo_details,
            **pillow_details,
        }
    return image_path, {
        "sanitized": True,
        "external_resources_blocked": True,
        **cairo_details,
        **pillow_details,
    }


def image_is_nonblank(path: Path) -> tuple[bool, dict[str, Any]]:
    with Image.open(path).convert("RGB") as image:
        extrema = image.getextrema()
        diff = ImageChops.difference(image, Image.new("RGB", image.size, (255, 255, 255)))
        bbox = diff.getbbox()
        changed_ratio = 0.0
        if bbox:
            gray = diff.convert("L")
            histogram = gray.histogram()
            changed = sum(count for value, count in enumerate(histogram) if value > 8)
            changed_ratio = changed / (image.width * image.height)
        return bool(bbox and changed_ratio > 0.01), {
            "size": image.size,
            "extrema": extrema,
            "changed_ratio": round(changed_ratio, 4),
        }


def render_charts(report: Report, out_dir: Path, preflight: Preflight) -> None:
    for chart in report.charts:
        path, details = render_svg_chart(chart, out_dir)
        if path is None:
            preflight.add(
                f"chart_render:{chart.index}",
                False,
                f"chart {chart.index} could not be rendered",
                details,
            )
            continue
        try:
            ok, image_details = image_is_nonblank(path)
        except Exception as exc:
            preflight.add(
                f"chart_render:{chart.index}",
                False,
                f"chart {chart.index} rendered unreadable image",
                {"path": str(path), "error": str(exc)[:500], **details},
            )
            continue
        chart.image_path = str(path)
        preflight.add(
            f"chart_render:{chart.index}",
            ok,
            f"chart {chart.index} rendered to {path.name}"
            if ok
            else f"chart {chart.index} rendered blank",
            {**image_details, **details},
        )


def table_slide_title(table: TableItem) -> str:
    title = collapse_ws(table.title)
    if title:
        return title if "table" in normalized_label(title) else f"{title} table"
    if table.headers:
        return "Report table: " + ", ".join(table.headers[:3])
    return "Report table"


def conclusions_slide_title(report: Report) -> str:
    has_conclusions = has_section_content(report, CONCLUSION_SECTION_CANDIDATES)
    has_implications = has_section_content(report, IMPLICATION_SECTION_CANDIDATES)
    if has_conclusions and has_implications:
        return "Conclusions and implications"
    if has_conclusions:
        return section_title_by_candidates(report, CONCLUSION_SECTION_CANDIDATES, "Conclusions")
    if has_implications:
        return section_title_by_candidates(report, IMPLICATION_SECTION_CANDIDATES, "Implications")
    return "Key takeaways"


def slide_plan(report: Report) -> list[dict[str, Any]]:
    executive = section_items_by_candidates(report, EXECUTIVE_SECTION_CANDIDATES)
    conclusions = section_items_by_candidates(report, CONCLUSION_SECTION_CANDIDATES)
    implications = section_items_by_candidates(report, IMPLICATION_SECTION_CANDIDATES)
    caveats = section_items_by_candidates(report, CAVEAT_SECTION_CANDIDATES)
    sources_section = find_section(report, SOURCE_SECTION_CANDIDATES)
    sources = (sources_section.paragraphs + sources_section.bullets) if sources_section else []
    plan: list[dict[str, Any]] = [
        {
            "kind": "cover",
            "title": report.title,
            "elements": ["slide-title", "lede", "metric-1", "metric-2", "metric-3", "metric-4"],
        },
        {
            "kind": "executive_summary",
            "title": "Executive summary",
            "elements": [
                "slide-title",
                "summary-bullets",
                "metric-card-1",
                "metric-card-2",
                "metric-card-3",
                "metric-card-4",
            ],
            "source_paragraphs": executive,
        },
    ]
    for chart in report.charts:
        plan.append(
            {
                "kind": "chart_evidence",
                "title": chart.takeaway,
                "chart_index": chart.index,
                "elements": ["slide-title", "chart-main", "chart-callout", "source-note"],
            }
        )
    for table_index, table in enumerate(report.tables, start=1):
        plan.append(
            {
                "kind": "table",
                "title": table_slide_title(table),
                "table_index": table_index,
                "elements": ["slide-title", "table-main", "table-note", "source-note"],
            }
        )
    if conclusions or implications:
        plan.append(
            {
                "kind": "conclusions_implications",
                "title": conclusions_slide_title(report),
                "elements": ["slide-title", "conclusions", "implications"],
                "source_paragraphs": {"conclusions": conclusions, "implications": implications},
            }
        )
    if caveats:
        plan.append(
            {
                "kind": "caveats",
                "title": section_title_by_candidates(
                    report, CAVEAT_SECTION_CANDIDATES, "Caveats and assumptions"
                ),
                "elements": ["slide-title", "caveats"],
                "source_paragraphs": caveats,
            }
        )
    if sources_section and (sources or sources_section.links):
        source_bullets, source_links = source_bullets_with_links(report)
        source_title = section_title_by_candidates(
            report, SOURCE_SECTION_CANDIDATES, "Sources and reproducibility"
        )
        plan.append(
            {
                "kind": "sources_reproducibility",
                "title": source_title,
                "elements": ["slide-title", "sources"],
                "source_paragraphs": sources,
                "source_links": [
                    {"text": bullet, "href": link}
                    for bullet, link in zip(source_bullets, source_links)
                    if link
                ],
            }
        )
    return plan


def add_rule(slide: Any, name: str, x: float, y: float, w: float, color: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.025)
    )
    shape.name = name
    set_shape_fill(shape, color)
    set_shape_line(shape, color, width_pt=0)


def add_vertical_rule(slide: Any, name: str, x: float, y: float, h: float, color: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.015), Inches(h)
    )
    shape.name = name
    set_shape_fill(shape, color)
    set_shape_line(shape, color, width_pt=0)


def add_metric_card(
    slide: Any,
    metric: MetricCard,
    idx: int,
    x: float,
    y: float,
    w: float,
    h: float,
    colors: dict[str, str],
    preflight: Preflight,
) -> None:
    add_box(
        slide,
        f"metric-{idx}",
        x,
        y,
        w,
        h,
        "#FFFFFF",
        colors.get("grid", DEFAULT_COLORS["grid"]),
        radius=True,
    )
    add_text(
        slide,
        f"metric-label-{idx}",
        x + 0.12,
        y + 0.12,
        w - 0.24,
        0.32,
        metric.label,
        7.5,
        colors["muted"],
        preflight,
        bold=True,
        min_font=6,
    )
    add_text(
        slide,
        f"metric-value-{idx}",
        x + 0.12,
        y + 0.46,
        w - 0.24,
        0.46,
        metric.value,
        20,
        colors["blue"],
        preflight,
        bold=True,
        min_font=14,
    )
    add_text(
        slide,
        f"metric-note-{idx}",
        x + 0.12,
        y + 0.85,
        w - 0.24,
        0.36,
        metric.note,
        8.5,
        colors["muted"],
        preflight,
        min_font=6.5,
    )


def executive_bullets(report: Report) -> list[str]:
    items = section_items_by_candidates(report, EXECUTIVE_SECTION_CANDIDATES)
    if not items:
        return [report.lede]
    bullets = [first_sentence(item, 320) for item in items[:4]]
    return [bullet for bullet in bullets if bullet]


def shorten_for_bullet(text: str, max_chars: int) -> str:
    text = collapse_ws(text)
    if len(text) <= max_chars:
        return text
    sentences = re.split(r"(?<=[.!?])\s+", text)
    if sentences and len(sentences[0]) <= max_chars:
        return sentences[0]
    return text[: max_chars - 3].rsplit(" ", 1)[0] + "..."


def split_sentences(paragraphs: list[str], max_items: int = 6, max_chars: int = 170) -> list[str]:
    output: list[str] = []
    for paragraph in paragraphs:
        paragraph = collapse_ws(paragraph)
        if not paragraph:
            continue
        if len(paragraph) <= max_chars:
            output.append(paragraph)
            if len(output) >= max_items:
                return output
            continue
        parts = re.split(r"(?<=[.!?])\s+", paragraph)
        for part in parts:
            part = collapse_ws(part)
            if part:
                output.append(shorten_for_bullet(part, max_chars))
            if len(output) >= max_items:
                return output
    return output


def list_slide_items(items: list[str], max_items: int = 6, max_chars: int = 220) -> list[str]:
    cleaned = [collapse_ws(item) for item in items if collapse_ws(item)]
    if len(cleaned) <= max_items:
        return [shorten_for_bullet(item, max_chars) for item in cleaned]
    return split_sentences(cleaned, max_items=max_items, max_chars=max_chars)


def normalize_sentence_key(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", collapse_ws(text).lower()).strip()


def chart_callout_text(chart: ChartItem) -> str:
    text = collapse_ws(chart.caveat)
    if not text:
        return ""
    sentences = [
        collapse_ws(part) for part in re.split(r"(?<=[.!?])\s+", text) if collapse_ws(part)
    ]
    if not sentences:
        return text

    first_key = normalize_sentence_key(sentences[0])
    takeaway_key = normalize_sentence_key(chart.takeaway)
    title_key = normalize_sentence_key(chart.title)
    if (
        first_key
        and (takeaway_key or title_key)
        and (
            first_key == takeaway_key
            or first_key.startswith((takeaway_key, title_key))
            or takeaway_key.startswith(first_key)
            or first_key == title_key
            or title_key.startswith(first_key)
        )
    ):
        trimmed = " ".join(sentences[1:]).strip()
        if trimmed:
            return trimmed
        return text
    return text


def build_cover(prs: Presentation, report: Report, preflight: Preflight) -> None:
    colors = report.colors
    slide = blank_slide(prs, colors)
    add_text(
        slide,
        "eyebrow",
        0.72,
        0.46,
        5.2,
        0.36,
        report.eyebrow or "Analytics report",
        9.5,
        colors["muted"],
        preflight,
        bold=True,
        min_font=7,
    )
    add_text(
        slide,
        "slide-title",
        0.68,
        0.95,
        11.7,
        1.55,
        report.title,
        28,
        colors["ink"],
        preflight,
        bold=True,
        min_font=20,
    )
    add_rule(slide, "title-rule", 0.72, 2.7, 1.45, colors["gold"])
    add_text(
        slide,
        "lede",
        0.72,
        3.02,
        10.7,
        1.0,
        report.lede,
        14,
        colors["neutral_dark"],
        preflight,
        min_font=10,
    )
    metric_count = min(4, len(report.metrics))
    if metric_count:
        gap = 0.18
        card_w = (11.9 - gap * (metric_count - 1)) / metric_count
        for idx, metric in enumerate(report.metrics[:metric_count], start=1):
            add_metric_card(
                slide,
                metric,
                idx,
                0.72 + (idx - 1) * (card_w + gap),
                5.45,
                card_w,
                1.24,
                colors,
                preflight,
            )
    add_source(slide, "Source: HTML report conversion inventory.", colors, preflight)


def build_executive(prs: Presentation, report: Report, preflight: Preflight) -> None:
    colors = report.colors
    slide = blank_slide(prs, colors)
    add_text(
        slide,
        "slide-title",
        0.65,
        0.45,
        12.0,
        0.6,
        "Executive summary",
        25,
        colors["ink"],
        preflight,
        bold=True,
        min_font=18,
    )
    add_rule(slide, "summary-rule", 0.72, 1.17, 1.2, colors["gold"])
    if not report.metrics:
        add_bullets(
            slide,
            "summary-bullets",
            0.72,
            1.48,
            11.55,
            4.75,
            executive_bullets(report),
            16.2,
            colors["ink"],
            preflight,
            max_items=5,
            min_font=11,
        )
        add_source(slide, "Source: report executive summary.", colors, preflight)
        return
    add_bullets(
        slide,
        "summary-bullets",
        0.72,
        1.48,
        7.08,
        4.55,
        executive_bullets(report),
        15.7,
        colors["ink"],
        preflight,
        max_items=4,
        min_font=11,
    )
    add_text(
        slide,
        "metrics-heading",
        8.32,
        1.22,
        4.25,
        0.34,
        "Key metrics",
        10.5,
        colors["muted"],
        preflight,
        bold=True,
        min_font=8,
    )
    metric_w = 2.02
    metric_h = 1.38
    gap_x = 0.2
    gap_y = 0.26
    for idx, metric in enumerate(report.metrics[:4], start=1):
        col = (idx - 1) % 2
        row = (idx - 1) // 2
        x = 8.32 + col * (metric_w + gap_x)
        y = 1.68 + row * (metric_h + gap_y)
        add_box(
            slide,
            f"metric-card-{idx}",
            x,
            y,
            metric_w,
            metric_h,
            "#FFFFFF",
            colors["grid"],
            radius=True,
        )
        add_text(
            slide,
            f"summary-metric-label-{idx}",
            x + 0.14,
            y + 0.12,
            metric_w - 0.28,
            0.28,
            metric.label,
            8.4,
            colors["muted"],
            preflight,
            bold=True,
            min_font=7,
        )
        add_text(
            slide,
            f"summary-metric-value-{idx}",
            x + 0.14,
            y + 0.46,
            metric_w - 0.28,
            0.42,
            metric.value,
            19,
            colors["blue"],
            preflight,
            bold=True,
            min_font=14,
        )
        add_text(
            slide,
            f"summary-metric-note-{idx}",
            x + 0.14,
            y + 0.93,
            metric_w - 0.28,
            0.3,
            metric.note,
            8.8,
            colors["muted"],
            preflight,
            min_font=7,
        )
    add_source(slide, "Source: report executive summary and top metric cards.", colors, preflight)


def build_chart_slide(
    prs: Presentation, report: Report, chart: ChartItem, preflight: Preflight
) -> None:
    colors = report.colors
    slide = blank_slide(prs, colors)
    add_text(
        slide,
        "slide-title",
        0.62,
        0.38,
        12.0,
        0.82,
        chart.takeaway,
        20,
        colors["ink"],
        preflight,
        bold=True,
        min_font=13.5,
    )
    chart_box = add_box(
        slide, "chart-frame", 0.72, 1.28, 8.0, 5.05, "#FFFFFF", colors["grid"], radius=True
    )
    chart_box.shadow.inherit = False
    if not chart.image_path:
        preflight.add(f"chart_image:{chart.index}", False, "chart image path missing")
        return
    try:
        picture = slide.shapes.add_picture(
            chart.image_path, Inches(0.95), Inches(1.55), width=Inches(7.55)
        )
    except Exception as exc:
        preflight.add(
            f"chart_image:{chart.index}",
            False,
            "chart image could not be inserted into PPTX",
            {"path": chart.image_path, "error": str(exc)[:500]},
        )
        return
    picture.name = "chart-main"
    if picture.height > Inches(4.45):
        picture.height = Inches(4.45)
    picture.left = Inches(0.95 + (7.55 - picture.width / 914400) / 2)
    picture.top = Inches(1.55 + (4.45 - picture.height / 914400) / 2)
    add_box(
        slide, "chart-callout-frame", 9.0, 1.28, 3.55, 5.05, "#FFFFFF", colors["grid"], radius=True
    )
    callout = chart_callout_text(chart)
    add_text(
        slide,
        "chart-callout-heading",
        9.25,
        1.55,
        3.05,
        0.35,
        "Context" if callout else "Source",
        10,
        colors["muted"],
        preflight,
        bold=True,
        min_font=7,
    )
    if callout:
        add_text(
            slide,
            "chart-callout",
            9.25,
            1.95,
            3.05,
            2.2,
            callout,
            12.5,
            colors["ink"],
            preflight,
            min_font=8.5,
        )
    add_text(
        slide,
        "chart-source",
        9.25,
        4.68 if callout else 1.95,
        3.05,
        0.75 if callout else 1.4,
        chart.title,
        8.5,
        colors["muted"],
        preflight,
        min_font=6.5,
    )
    add_source(
        slide,
        f"Source chart {chart.index}: extracted from report SVG and rendered as chart evidence.",
        colors,
        preflight,
    )


def build_table_slide(
    prs: Presentation, report: Report, table: TableItem, preflight: Preflight
) -> None:
    colors = report.colors
    slide = blank_slide(prs, colors)
    add_text(
        slide,
        "slide-title",
        0.62,
        0.42,
        12.0,
        0.55,
        table_slide_title(table),
        22,
        colors["ink"],
        preflight,
        bold=True,
        min_font=16,
    )
    cols = max([1, len(table.headers), *[len(row) for row in table.rows]])
    headers = table.headers + [f"Column {idx}" for idx in range(len(table.headers) + 1, cols + 1)]
    body_rows = [row + [""] * (cols - len(row)) for row in table.rows]
    rows = max(1, len(body_rows) + 1)
    table_h = min(4.25, max(3.05, 0.58 + max(1, len(table.rows)) * 0.42))
    shape = slide.shapes.add_table(
        rows, cols, Inches(0.62), Inches(1.35), Inches(12.1), Inches(table_h)
    )
    shape.name = "table-main"
    ppt_table = shape.table
    if cols >= 1:
        ppt_table.columns[0].width = Inches(1.95)
    for c in range(1, cols):
        ppt_table.columns[c].width = Inches((12.1 - 1.95) / max(1, cols - 1))
    ppt_table.rows[0].height = Inches(0.52)
    body_h = table_h - 0.52
    if rows > 1:
        body_h = (table_h - 0.52) / (rows - 1)
        for r in range(1, rows):
            ppt_table.rows[r].height = Inches(body_h)
    body_font = 11.4 if body_h >= 0.38 else 10.2 if body_h >= 0.31 else 9.2
    for c, header in enumerate(headers):
        cell = ppt_table.cell(0, c)
        cell.text = header
        set_shape_fill(cell, colors["blue_light"])
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(10.2)
                run.font.bold = True
                run.font.color.rgb = rgb(colors["ink"])
    for r, row in enumerate(body_rows, start=1):
        for c in range(cols):
            cell = ppt_table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            set_shape_fill(cell, "#FFFFFF")
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(body_font if c else body_font + 0.4)
                    run.font.bold = c == 0
                    run.font.color.rgb = rgb(colors["ink"])
    note = table.note or "Source table extracted from the report."
    add_text(
        slide,
        "table-note",
        0.78,
        1.35 + table_h + 0.32,
        11.7,
        0.55,
        note,
        11.5,
        colors["muted"],
        preflight,
        min_font=9,
    )
    preflight.add(
        "table_coverage:main",
        bool(headers and body_rows),
        f"main table has {len(body_rows)} body rows and {cols} columns",
        {"rows": len(body_rows), "columns": cols},
    )
    add_source(slide, "Source: report metric table.", colors, preflight)


def build_conclusions(prs: Presentation, report: Report, preflight: Preflight) -> None:
    colors = report.colors
    slide = blank_slide(prs, colors)
    add_text(
        slide,
        "slide-title",
        0.62,
        0.42,
        12.0,
        0.72,
        conclusions_slide_title(report),
        21,
        colors["ink"],
        preflight,
        bold=True,
        min_font=15,
    )
    conclusion_source_items = section_items_by_candidates(report, CONCLUSION_SECTION_CANDIDATES)
    implication_source_items = section_items_by_candidates(report, IMPLICATION_SECTION_CANDIDATES)
    conclusions = list_slide_items(
        conclusion_source_items,
        max_items=5 if not implication_source_items else 4,
        max_chars=320 if not implication_source_items else 220,
    )
    implications = list_slide_items(
        implication_source_items,
        max_items=6 if not conclusion_source_items else 5,
        max_chars=320 if not conclusion_source_items else 220,
    )
    conclusion_heading = section_title_by_candidates(
        report, CONCLUSION_SECTION_CANDIDATES, "Conclusions"
    )
    implication_heading = section_title_by_candidates(
        report, IMPLICATION_SECTION_CANDIDATES, "Implications"
    )
    if conclusions and not implications:
        add_bullets(
            slide,
            "conclusions",
            0.75,
            1.48,
            11.4,
            4.75,
            conclusions,
            14.6,
            colors["ink"],
            preflight,
            max_items=5,
            min_font=10.5,
        )
        add_source(slide, "Source: report conclusion section.", colors, preflight)
        return
    if implications and not conclusions:
        add_bullets(
            slide,
            "implications",
            0.75,
            1.48,
            11.4,
            4.75,
            implications,
            14.6,
            colors["ink"],
            preflight,
            max_items=5,
            min_font=10.5,
        )
        add_source(slide, "Source: report implications section.", colors, preflight)
        return
    add_text(
        slide,
        "conclusions-heading",
        0.78,
        1.38,
        5.2,
        0.35,
        conclusion_heading,
        13,
        colors["blue"],
        preflight,
        bold=True,
        min_font=10,
    )
    add_bullets(
        slide,
        "conclusions",
        0.75,
        1.82,
        5.55,
        4.35,
        conclusions,
        14.2,
        colors["ink"],
        preflight,
        max_items=4,
        min_font=10,
    )
    add_text(
        slide,
        "implications-heading",
        7.02,
        1.38,
        5.2,
        0.35,
        implication_heading,
        13,
        colors["blue"],
        preflight,
        bold=True,
        min_font=10,
    )
    add_bullets(
        slide,
        "implications",
        7.0,
        1.82,
        5.55,
        4.35,
        implications,
        14.2,
        colors["ink"],
        preflight,
        max_items=4,
        min_font=10,
    )
    add_vertical_rule(slide, "column-divider", 6.67, 1.35, 4.95, colors["grid"])
    add_source(slide, "Source: report conclusions and implications sections.", colors, preflight)


def source_bullets_with_links(report: Report) -> tuple[list[str], list[str | None]]:
    section = find_section(report, SOURCE_SECTION_CANDIDATES) or Section("Sources")
    bullets: list[str] = []
    links: list[str | None] = []
    for link in section.links:
        label = collapse_ws(link.text)
        href = collapse_ws(link.href)
        if not href:
            continue
        display = href if not label or label == href else f"{label}: {href}"
        bullets.append(display)
        links.append(href)

    remaining_text_slots = max(0, 6 - len(bullets))
    source_text = (
        split_sentences(
            section.paragraphs + section.bullets,
            max_items=remaining_text_slots,
            max_chars=300,
        )
        if remaining_text_slots
        else []
    )
    for item in source_text:
        if any(item in existing or existing in item for existing in bullets):
            continue
        bullets.append(item)
        links.append(None)
    return bullets, links


def build_caveats_slide(prs: Presentation, report: Report, preflight: Preflight) -> None:
    colors = report.colors
    slide = blank_slide(prs, colors)
    add_text(
        slide,
        "slide-title",
        0.62,
        0.42,
        12.0,
        0.55,
        section_title_by_candidates(report, CAVEAT_SECTION_CANDIDATES, "Caveats and assumptions"),
        23,
        colors["ink"],
        preflight,
        bold=True,
        min_font=17,
    )
    add_rule(slide, "caveats-rule", 0.72, 1.18, 1.2, colors["gold"])
    caveats = split_sentences(
        section_items_by_candidates(report, CAVEAT_SECTION_CANDIDATES),
        max_items=5,
        max_chars=205,
    )
    add_bullets(
        slide,
        "caveats",
        0.92,
        1.52,
        11.15,
        4.85,
        caveats,
        15.1,
        colors["ink"],
        preflight,
        max_items=5,
        min_font=11,
    )
    add_source(slide, "Source: report caveats and assumptions section.", colors, preflight)


def build_sources_slide(
    prs: Presentation,
    report: Report,
    preflight: Preflight,
    bullets: list[str],
    links: list[str | None],
) -> None:
    colors = report.colors
    slide = blank_slide(prs, colors)
    base_title = section_title_by_candidates(
        report, SOURCE_SECTION_CANDIDATES, "Sources and reproducibility"
    )
    add_text(
        slide,
        "slide-title",
        0.62,
        0.42,
        12.0,
        0.55,
        base_title,
        23,
        colors["ink"],
        preflight,
        bold=True,
        min_font=17,
    )
    add_rule(slide, "sources-rule", 0.72, 1.18, 1.2, colors["gold"])
    add_bullets(
        slide,
        "sources",
        0.92,
        1.38,
        11.15,
        5.45,
        bullets,
        10.6,
        colors["ink"],
        preflight,
        max_items=len(bullets),
        min_font=6,
        links=links,
        link_color=colors["blue"],
        space_after_pt=2.0,
    )
    all_links = [link for link in source_bullets_with_links(report)[1] if link]
    preflight.add(
        "coverage:source_links",
        True,
        f"preserved {len(all_links)} source links from the report",
        {"links": all_links},
    )
    add_source(slide, "Source: report sources and reproducibility section.", colors, preflight)


def create_pptx(report: Report, out_dir: Path, preflight: Preflight) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_cover(prs, report, preflight)
    build_executive(prs, report, preflight)
    for chart in report.charts:
        build_chart_slide(prs, report, chart, preflight)
    for table in report.tables:
        build_table_slide(prs, report, table, preflight)
    if has_section_content(report, CONCLUSION_SECTION_CANDIDATES) or has_section_content(
        report, IMPLICATION_SECTION_CANDIDATES
    ):
        build_conclusions(prs, report, preflight)
    if has_section_content(report, CAVEAT_SECTION_CANDIDATES):
        build_caveats_slide(prs, report, preflight)
    if has_section_content(report, SOURCE_SECTION_CANDIDATES):
        source_bullets, source_links = source_bullets_with_links(report)
        build_sources_slide(prs, report, preflight, source_bullets, source_links)
    pptx_path = out_dir / "deck.pptx"
    prs.save(pptx_path)
    return pptx_path


def pptx_package_ok(path: Path) -> tuple[bool, dict[str, Any]]:
    try:
        with zipfile.ZipFile(path) as archive:
            bad = archive.testzip()
            slide_files = [
                name for name in archive.namelist() if re.match(r"ppt/slides/slide\d+\.xml", name)
            ]
            media_files = [name for name in archive.namelist() if name.startswith("ppt/media/")]
            return bad is None and bool(slide_files), {
                "bad_member": bad,
                "slide_count": len(slide_files),
                "media_count": len(media_files),
            }
    except zipfile.BadZipFile:
        return False, {"error": "bad zip file"}


def build_manifest(report: Report) -> dict[str, Any]:
    return {
        "source": report.source,
        "title": report.title,
        "eyebrow": report.eyebrow,
        "counts": {
            "metrics": len(report.metrics),
            "sections": len(report.sections),
            "charts": len(report.charts),
            "tables": len(report.tables),
        },
        "metrics": [asdict(metric) for metric in report.metrics],
        "sections": {title: asdict(section) for title, section in report.sections.items()},
        "charts": [
            {
                "index": chart.index,
                "title": chart.title,
                "takeaway": chart.takeaway,
                "caveat": chart.caveat,
                "image_path": chart.image_path,
                "width_px": chart.width_px,
                "height_px": chart.height_px,
            }
            for chart in report.charts
        ],
        "tables": [asdict(table) for table in report.tables],
        "colors": report.colors,
    }


def add_coverage_checks(report: Report, preflight: Preflight) -> None:
    preflight.add(
        "source:title", bool(report.title), "source report title extracted", {"title": report.title}
    )
    has_content = bool(
        report.sections or report.metrics or report.charts or report.tables or report.lede
    )
    preflight.add(
        "source:content",
        has_content,
        "source report body content extracted",
        {
            "sections": len(report.sections),
            "metrics": len(report.metrics),
            "charts": len(report.charts),
            "tables": len(report.tables),
        },
    )
    preflight.add(
        "source:metrics",
        True,
        f"extracted {len(report.metrics)} metric cards",
        {"metrics": len(report.metrics)},
    )
    preflight.add(
        "source:charts",
        True,
        f"extracted {len(report.charts)} charts",
        {"charts": len(report.charts)},
    )
    preflight.add(
        "source:tables",
        True,
        f"extracted {len(report.tables)} tables",
        {"tables": len(report.tables)},
    )
    preflight.add(
        "coverage:executive_section",
        True,
        "executive-style section discovery completed",
        {"found": bool(find_section(report, EXECUTIVE_SECTION_CANDIDATES))},
    )
    preflight.add(
        "coverage:caveats",
        True,
        "caveat-style section discovery completed",
        {"found": bool(find_section(report, CAVEAT_SECTION_CANDIDATES))},
    )
    preflight.add(
        "coverage:sources",
        True,
        "source-style section discovery completed",
        {"found": bool(find_section(report, SOURCE_SECTION_CANDIDATES))},
    )


def run(source: str, out_dir: Path) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    preflight = Preflight()
    source_path = resolve_source(source)
    report = parse_report(source_path)
    add_coverage_checks(report, preflight)
    render_charts(report, out_dir, preflight)
    plan = slide_plan(report)
    write_json(out_dir / "deck_plan.json", plan)
    pptx_path = create_pptx(report, out_dir, preflight)
    ok, details = pptx_package_ok(pptx_path)
    preflight.add("pptx:package", ok, "PPTX package is readable", details)
    manifest = build_manifest(report)
    manifest["outputs"] = {
        "deck_pptx": str(pptx_path),
        "deck_plan": str(out_dir / "deck_plan.json"),
        "preflight": str(out_dir / "preflight_checks.json"),
    }
    write_json(out_dir / "manifest.json", manifest)
    write_json(out_dir / "preflight_checks.json", preflight.to_json())
    print(
        json.dumps(
            {
                "status": "passed" if preflight.passed else "failed",
                "pptx": str(pptx_path),
                "out_dir": str(out_dir),
            },
            indent=2,
        )
    )
    return 0 if preflight.passed else 2


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert an HTML analytics report into a local PPTX deck."
    )
    parser.add_argument("source", help="Local path to the HTML report file.")
    parser.add_argument(
        "--out-dir", required=True, help="Directory for deck.pptx and verification artifacts."
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv or sys.argv[1:])
    return run(args.source, Path(args.out_dir).expanduser().resolve())


if __name__ == "__main__":
    raise SystemExit(main())
